import PyPDF2
from pptx import Presentation

def pdf_to_text(pdf_file):
    # Extract text from a PDF file
    with open(pdf_file, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text_pages = []
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text_pages.append(page.extract_text())
    return text_pages

def create_ppt(text_pages, ppt_file):
    prs = Presentation()

    for i, text in enumerate(text_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]

        title.text = f"Page {i+1}"
        content.text = text if text else "(No text extracted)"

    prs.save(ppt_file)
    print(f"PPT saved as: {ppt_file}")

# >>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
# USE YOUR PDF FILE HERE
pdf_path = r"C:\Users\jsrin\OneDrive\Desktop\PEKG\MTP_PRESENTATION (1).pdf"
ppt_output = r"C:\Users\jsrin\OneDrive\Desktop\PEKG\converted_presentation.pptx"
# >>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>

pages = pdf_to_text(pdf_path)
create_ppt(pages, ppt_output)
